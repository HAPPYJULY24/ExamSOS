# modules/file_handler.py
# 文件上传/预览 & PDF 导出（Markdown → PDF，带格式，稳定版）

import streamlit as st
import os
from PyPDF2 import PdfReader
import docx
import markdown2
from bs4 import BeautifulSoup
import pdfkit
from pptx import Presentation  # ✅ 新增：支持 PowerPoint


# 预览上传的文件
def preview_files(uploaded_files):
    for file in uploaded_files:
        st.write(f"📄 文件名: {file.name} ({file.size / 1024:.2f} KB)")


# 读取 PPTX 内容
def read_pptx_content(uploaded_file):
    prs = Presentation(uploaded_file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_runs.append(shape.text.strip())
    return "\n".join(text_runs)


# 读取文件内容（支持 PDF、DOCX、TXT、PPTX）
def read_file_content(uploaded_file):
    if uploaded_file.type == "application/pdf":
        pdf_reader = PdfReader(uploaded_file)
        text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
        return text

    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = docx.Document(uploaded_file)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return text

    elif uploaded_file.type == "text/plain":
        return uploaded_file.read().decode("utf-8", errors="ignore")

    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return read_pptx_content(uploaded_file)  # ✅ 新增 PPTX 支持

    else:
        return "⚠️ 暂不支持该文件格式"


# 保存 AI 输出为 PDF (Markdown → PDF，带格式)
def save_to_pdf(markdown_text, filename="output.pdf", spacing=1):
    import re
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    output_path = os.path.join("outputs", filename)
    os.makedirs("outputs", exist_ok=True)

    # 清理非法字符（避免 PDF 崩溃）
    clean_text = re.sub(r"[^\x09\x0A\x0D\x20-\x7E\u4e00-\u9fff]", "", markdown_text)

    # Markdown → HTML
    text = markdown2.markdown(clean_text, extras=["fenced-code-blocks", "tables"])
    soup = BeautifulSoup(text, "html.parser")
    html_text = str(soup)

    # 包裹 CSS
    html_with_style = f"""
    <html>
    <head>
        <meta charset="utf-8">
        <style>
            body {{
                font-family: "Times New Roman", Times, serif;
                font-size: 12pt;
                line-height: 1.6;
                margin: 50px;
                letter-spacing: {spacing}px;
            }}
            h1 {{ font-size: 20pt; font-weight: bold; text-align: center; }}
            h2 {{ font-size: 16pt; font-weight: bold; margin-top: 18px; }}
            h3 {{ font-size: 14pt; font-weight: bold; margin-top: 14px; }}
            p, li {{ font-size: 12pt; margin-bottom: 6px; text-align: justify; }}
            pre, code {{
                background-color: #f5f5f5;
                border: 1px solid #ddd;
                padding: 6px;
                display: block;
                font-family: "Courier New", monospace;
                font-size: 10pt;
                white-space: pre-wrap;
                border-radius: 4px;
                margin: 10px 0;
            }}
            .math-block {{
                border: 1px solid #000;
                padding: 8px;
                text-align: center;
                margin: 12px 0;
                font-family: "Times New Roman", serif;
                font-style: italic;
                background-color: #fafafa;
            }}
            table {{
                border-collapse: collapse;
                width: 100%;
                margin: 12px 0;
            }}
            th, td {{
                border: 1px solid #000;
                padding: 6px;
                text-align: center;
                font-size: 10pt;
            }}
            th {{
                background-color: #f0f0f0;
                font-weight: bold;
            }}
            mark {{
                background-color: #ffeb3b;
                font-weight: bold;
                padding: 2px 4px;
                border-radius: 2px;
            }}
        </style>
    </head>
    <body>
        {html_text}
    </body>
    </html>
    """

    # wkhtmltopdf 配置
    path_wkhtmltopdf = r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe"
    config = pdfkit.configuration(wkhtmltopdf=path_wkhtmltopdf)

    options = {
        "enable-local-file-access": None,
        "encoding": "UTF-8",
        "quiet": "",
        "disable-smart-shrinking": "",
        "load-error-handling": "ignore",
        "load-media-error-handling": "ignore",
        "zoom": "1.0"
    }

    try:
        pdfkit.from_string(html_with_style, output_path, configuration=config, options=options)
    except Exception:
        # fallback → ReportLab 简单文本导出
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        y = height - 50
        c.setFont("Times-Roman", 12)
        for line in clean_text.split("\n"):
            c.drawString(50, y, line.strip())
            y -= 16
            if y < 50:  # 翻页
                c.showPage()
                c.setFont("Times-Roman", 12)
                y = height - 50
        c.save()

    return output_path
