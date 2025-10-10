# modules/file_parser.py
# 文件格式的转换，文件内容的提取
# 输入处理器（上传 → 提取文字）

import os
import io
import tempfile
from pptx import Presentation
from docx import Document
import fitz  # PyMuPDF
from PIL import Image, ImageStat
import numpy as np
import streamlit as st
from modules.logger import log_event   # ✅ 引入日志模块

# === 模块健康状态上报 ===
from modules.utils.system_status import update_module_status


try:
    update_module_status("file_parser", "active", "模块加载成功")
except Exception as e:
    update_module_status("file_parser", "error", str(e))



# ==============================
# OCR 功能占位（不加载 easyocr）
# ==============================

def ocr_image(pil_img):
    """OCR 识别图片 -> 文字（目前禁用，直接返回空字符串）"""
    return ""


def is_text_image(pil_img, threshold=0.05):
    """简单判断图片是否可能包含文字（逻辑保留）"""
    gray = pil_img.convert("L")
    stat = ImageStat.Stat(gray)
    variance = stat.var[0] / 255**2
    return variance > threshold


# ======================================================
#   各种文件格式解析
# ======================================================

def extract_text_from_pptx_file(file_bytes, filename="unknown.pptx"):
    """提取 PPTX 文本（OCR 暂时禁用）"""
    try:
        log_event("file_parser", "INFO", "work", f"开始解析 PPTX 文件: {filename}")
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, filename)
            with open(pptx_path, "wb") as f:
                f.write(file_bytes)

            prs = Presentation(pptx_path)
            text = []

            for i, slide in enumerate(prs.slides, start=1):
                slide_text = []

                # 提取可编辑文字
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            para_text = para.text.strip()
                            if para_text:
                                slide_text.append(para_text)

                # 图片 OCR（禁用）
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        try:
                            image = shape.image
                            pil_img = Image.open(io.BytesIO(image.blob))
                            if is_text_image(pil_img):
                                ocr_text = ocr_image(pil_img)
                                if ocr_text:
                                    slide_text.append(ocr_text)
                        except Exception as e:
                            log_event("file_parser", "WARNING", "warning", f"PPTX 图片OCR失败: {filename}", remark=str(e))

                if slide_text:
                    text.append(f"【Slide {i} - {filename}】\n" + "\n".join(set(slide_text)))

        result = "\n\n".join(text) if text else "（未提取到有效文本）"
        log_event("file_parser", "INFO", "work", f"PPTX 文件解析完成: {filename}", meta={"text_length": len(result)})
        return result

    except Exception as e:
        log_event("file_parser", "ERROR", "down", f"PPTX 解析失败: {filename}", remark=str(e))
        return f"❌ 文件解析失败: {filename}"


def extract_text_from_file(uploaded_file, filename=None):
    """根据文件类型提取纯文本（支持 PPTX/DOCX/PDF/TXT，OCR 暂时禁用）"""
    filename = (filename or getattr(uploaded_file, "name", "unknown")).lower()
    try:
        log_event("file_parser", "INFO", "work", f"开始解析文件: {filename}")

        # 确保 bytes 不会被读空
        if hasattr(uploaded_file, "seek"):
            uploaded_file.seek(0)
        file_bytes = uploaded_file.read() if hasattr(uploaded_file, "read") else uploaded_file
        file_stream = io.BytesIO(file_bytes)

        # -------- PPTX --------
        if filename.endswith(".pptx"):
            return extract_text_from_pptx_file(file_bytes, filename=filename)

        # -------- DOCX --------
        elif filename.endswith(".docx"):
            doc = Document(file_stream)
            text = []

            for p in doc.paragraphs:
                if p.text.strip():
                    text.append(p.text)

            # 图片 OCR（禁用）
            for rel in doc.part.rels.values():
                if hasattr(rel, "target_ref") and "image" in rel.target_ref:
                    try:
                        image_data = rel.target_part.blob
                        pil_img = Image.open(io.BytesIO(image_data))
                        ocr_text = ocr_image(pil_img)
                        if ocr_text:
                            text.append(ocr_text)
                    except Exception as e:
                        log_event("file_parser", "WARNING", "warning", f"DOCX 图片OCR失败: {filename}", remark=str(e))

            result = "\n".join(text) if text else "（未提取到有效文本）"
            log_event("file_parser", "INFO", "work", f"DOCX 文件解析完成: {filename}", meta={"text_length": len(result)})
            return result

        # -------- PDF --------
        elif filename.endswith(".pdf"):
            text = []
            with fitz.open(stream=file_bytes, filetype="pdf") as pdf_doc:
                for page_num, page in enumerate(pdf_doc, start=1):
                    page_text = page.get_text("text").strip()
                    ocr_text = ""  # OCR 暂不启用
                    combined_parts = []
                    if page_text:
                        combined_parts.append(page_text)
                    if ocr_text:
                        combined_parts.append(ocr_text)

                    page_content = "\n".join(combined_parts).strip()
                    text.append(f"【第 {page_num} 页 - {filename}】\n{page_content}")

            result = "\n\n".join(text) if text else "（未提取到有效文本）"
            log_event("file_parser", "INFO", "work", f"PDF 文件解析完成: {filename}", meta={"pages": len(text)})
            return result

        # -------- TXT --------
        elif filename.endswith(".txt"):
            try:
                result = file_bytes.decode("utf-8", errors="ignore").strip() or "（未提取到有效文本）"
                log_event("file_parser", "INFO", "work", f"TXT 文件解析完成: {filename}", meta={"text_length": len(result)})
                return result
            except Exception as e:
                log_event("file_parser", "ERROR", "down", f"TXT 文件解析失败: {filename}", remark=str(e))
                return "（未提取到有效文本）"

        else:
            log_event("file_parser", "WARNING", "warning", f"不支持的文件格式: {filename}")
            return f"❌ 不支持的文件格式: {filename}"

    except Exception as e:
        log_event("file_parser", "ERROR", "down", f"文件解析失败: {filename}", remark=str(e))
        return f"❌ 文件解析失败: {filename}"


def preview_files(uploaded_files):
    """只展示文件名"""
    try:
        log_event("file_parser", "INFO", "work", "预览上传文件列表", meta={"count": len(uploaded_files)})
        for uf in uploaded_files:
            st.write(f"✅ {uf.name}")
    except Exception as e:
        log_event("file_parser", "ERROR", "down", "预览文件失败", remark=str(e))


def merge_files_text(uploaded_files):
    """把多个文件内容拼接成一个大字符串，带文件名分隔符"""
    all_texts = []
    try:
        log_event("file_parser", "INFO", "work", "开始合并文件内容", meta={"count": len(uploaded_files)})
        for idx, uploaded_file in enumerate(uploaded_files, start=1):
            uploaded_file.seek(0)
            text = extract_text_from_file(uploaded_file)
            print(f"======= {uploaded_file.name} 内容预览 =======")
            print(text[:500])  # 限制调试输出，避免日志过大

            labeled_text = f"\n======= 文件 {idx}: {uploaded_file.name} =======\n{text}\n"
            all_texts.append(labeled_text)

        merged_text = "\n".join(all_texts)
        log_event("file_parser", "INFO", "work", "文件合并完成", meta={"total_length": len(merged_text)})
        return merged_text

    except Exception as e:
        log_event("file_parser", "ERROR", "down", "合并文件失败", remark=str(e))
        return "❌ 文件合并失败"
